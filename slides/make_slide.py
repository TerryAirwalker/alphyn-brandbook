#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
РЕФЕРЕНС-ШАБЛОН редактируемого слайда Alphyn (тип «Данные», Режим A).
Проверен на манифесте v1.1: крупная цифра + вывод + две icon_card с lucide-иконками.

Как работает (Режим A по manifest.slides.editability):
  1. Декор (aurora_ink + blueprint-сетка + контейнеры карточек + lucide-иконки + знак)
     рендерится в full-bleed PNG через Playwright.
  2. Весь ТЕКСТ — живые надписи python-pptx поверх картинки (редактируемо).
  3. Шрифты Onest/JetBrains Mono встраиваются отдельным шагом: embed_fonts.py.

Иконки/цвета/типографика берутся из брендбука (assets/icons/icons.json + токены манифеста).
Адаптируй тексты/цифры под задачу — цифры только из products/alphyn.ai, не выдумывать.

Запуск:  python3 make_slide.py   (-> out/test_slide.pptx, out/slide_full.png)
         python3 embed_fonts.py out/test_slide.pptx
"""
import json, base64, pathlib

HERE = pathlib.Path(__file__).parent
ROOT = HERE / "out"; ROOT.mkdir(exist_ok=True)
ASSETS = (HERE / ".." / "assets").resolve()
ICONS = json.load(open(ASSETS / "icons/icons.json"))
FONT_DIR = ASSETS / "fonts-embed"

# ---- токены (core, из манифеста) ----
INK = "#08060F"; INK2 = "#0E0A1C"
VIOLET = "#6E40D7"; VBRIGHT = "#8A5BF0"; INDIGO = "#22186E"; LAV = "#B8ABF4"
TEXT = "#EDEAFB"; SOFT = "#C9C2E6"; MUTED = "#9D95C4"
LINE = "rgba(184,171,244,.10)"
NBSP = " "

W, H = 1920, 1080
EMU = 6350  # EMU на px (12192000/1920)
CARD = dict(x=1170, w=620, h=205)
CARD1_Y, CARD2_Y = 318, 563
ICON_BOX, INSET = 56, 30


def icon_svg(name, color=VBRIGHT, box=ICON_BOX, gl=26):
    pad = (box - gl) / 2
    return (f'<svg width="{gl}" height="{gl}" viewBox="0 0 24 24" fill="none" stroke="{color}" '
            f'stroke-width="1.75" stroke-linecap="round" stroke-linejoin="round" '
            f'style="position:absolute;left:{pad}px;top:{pad}px">{ICONS[name]}</svg>')


def mark_glyph(h=24):
    b64 = base64.b64encode((ASSETS / "logo/mark.svg").read_text().encode()).decode()
    return f'<img src="data:image/svg+xml;base64,{b64}" style="height:{h}px;display:block">'


def card_html(cy, icon):
    c = CARD
    return f'''<div style="position:absolute;left:{c['x']}px;top:{cy}px;width:{c['w']}px;height:{c['h']}px;
      background:{INK2};border:1px solid {LINE};border-radius:18px;box-sizing:border-box">
      <div style="position:absolute;left:{INSET}px;top:{INSET}px;width:{ICON_BOX}px;height:{ICON_BOX}px;
        background:rgba(138,91,240,.12);border-radius:14px">{icon_svg(icon)}</div></div>'''


def build_html(decor=True):
    fonts = ""
    if not decor:
        def face(fam, path, wght):
            b64 = base64.b64encode((FONT_DIR / path).read_bytes()).decode()
            return (f"@font-face{{font-family:'{fam}';font-weight:{wght};"
                    f"src:url(data:font/ttf;base64,{b64}) format('truetype')}}")
        fonts = ("<style>" + face("Onest", "OnestEmbedRegular.ttf", 400)
                 + face("Onest", "OnestEmbedBold.ttf", 800)
                 + face("JetBrains Mono", "JBMonoEmbedRegular.ttf", 400) + "</style>")
    c = CARD
    text_layer = ""
    if not decor:
        text_layer = f'''
        <div style="position:absolute;left:130px;top:300px;font-family:'JetBrains Mono';font-size:15px;
          letter-spacing:.18em;color:{LAV};text-transform:uppercase">ПРОДАКШЕН{NBSP}· ОТКРЫТЫЙ LAKEHOUSE</div>
        <div style="position:absolute;left:126px;top:330px;font-family:'Onest';font-weight:800;font-size:200px;
          letter-spacing:-.03em;line-height:.98;color:{TEXT}">3{NBSP}ПБ<span style="color:{VBRIGHT}">+</span></div>
        <div style="position:absolute;left:132px;top:600px;width:960px;font-family:'Onest';font-weight:400;
          font-size:31px;line-height:1.5;color:{SOFT}">данных под{NBSP}управлением одной платформы{NBSP}— на{NBSP}открытом
          lakehouse-контуре, без{NBSP}копий между системами и{NBSP}вендорского замка.</div>
        <div style="position:absolute;left:{c['x']+104}px;top:{CARD1_Y+30}px;font-family:'Onest';font-weight:700;
          font-size:25px;color:{TEXT}">Lakehouse</div>
        <div style="position:absolute;left:{c['x']+104}px;top:{CARD1_Y+72}px;width:{c['w']-134}px;font-family:'Onest';
          font-weight:400;font-size:17px;line-height:1.5;color:{MUTED}">StarRocks · Impala · Spark · Trino на{NBSP}Iceberg/S3</div>
        <div style="position:absolute;left:{c['x']+104}px;top:{CARD2_Y+30}px;font-family:'Onest';font-weight:700;
          font-size:25px;color:{TEXT}">AI Studio</div>
        <div style="position:absolute;left:{c['x']+104}px;top:{CARD2_Y+72}px;width:{c['w']-134}px;font-family:'Onest';
          font-weight:400;font-size:17px;line-height:1.5;color:{MUTED}">ML · GenAI · агентный AI, MLOps/LLMOps</div>
        <div style="position:absolute;left:168px;top:988px;font-family:'JetBrains Mono';font-size:16px;
          color:{MUTED};letter-spacing:.04em">alphyn.ai</div>'''
    return f'''<!doctype html><html><head><meta charset="utf-8">{fonts}
    <style>*{{margin:0;box-sizing:border-box}}html,body{{width:{W}px;height:{H}px;overflow:hidden}}</style></head>
    <body style="width:{W}px;height:{H}px;position:relative;background:{INK}">
      <div style="position:absolute;inset:0;background-image:linear-gradient({LINE} 1px,transparent 1px),
        linear-gradient(90deg,{LINE} 1px,transparent 1px);background-size:64px 64px;
        -webkit-mask-image:radial-gradient(120% 90% at 70% 0%,#000 30%,transparent 78%)"></div>
      <div style="position:absolute;right:-120px;top:-160px;width:820px;height:820px;border-radius:50%;
        background:radial-gradient(circle,rgba(110,64,215,.55),rgba(34,24,110,.20) 45%,transparent 70%);filter:blur(64px)"></div>
      <div style="position:absolute;left:60px;bottom:-220px;width:620px;height:620px;border-radius:50%;
        background:radial-gradient(circle,rgba(138,91,240,.16),transparent 70%);filter:blur(70px)"></div>
      {card_html(CARD1_Y, "database")}
      {card_html(CARD2_Y, "brain")}
      <div style="position:absolute;left:130px;top:988px">{mark_glyph(24)}</div>
      {text_layer}
    </body></html>'''


def render(html_path, png_path):
    from playwright.sync_api import sync_playwright
    with sync_playwright() as pw:
        b = pw.chromium.launch()
        p = b.new_page(viewport={"width": W, "height": H}, device_scale_factor=2)
        p.goto(pathlib.Path(html_path).as_uri(), wait_until="networkidle")
        p.wait_for_timeout(400)
        p.screenshot(path=str(png_path))
        b.close()


def main():
    (ROOT / "slide_decor.html").write_text(build_html(True))
    (ROOT / "slide_full.html").write_text(build_html(False))
    render(ROOT / "slide_decor.html", ROOT / "slide_decor.png")
    render(ROOT / "slide_full.html", ROOT / "slide_full.png")

    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.lang import MSO_LANGUAGE_ID

    prs = Presentation()
    prs.slide_width = Emu(W * EMU); prs.slide_height = Emu(H * EMU)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(ROOT / "slide_decor.png"), 0, 0, Emu(W * EMU), Emu(H * EMU))

    def hexrgb(h): return RGBColor.from_string(h.lstrip("#"))

    def tb(x, y, w, h, runs, size, weight, font="Onest", lh=1.0):
        box = slide.shapes.add_textbox(Emu(x * EMU), Emu(y * EMU), Emu(w * EMU), Emu(h * EMU))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        if lh: p.line_spacing = lh
        for txt, col in runs:
            r = p.add_run(); r.text = txt
            f = r.font; f.name = font; f.size = Pt(size); f.bold = (weight >= 700)
            f.color.rgb = hexrgb(col); f.language_id = MSO_LANGUAGE_ID.RUSSIAN

    c = CARD
    tb(130, 300, 900, 34, [("ПРОДАКШЕН" + NBSP + "· ОТКРЫТЫЙ LAKEHOUSE", LAV)], 11.5, 400, "JetBrains Mono")
    tb(126, 322, 1000, 250, [("3" + NBSP + "ПБ", TEXT), ("+", VBRIGHT)], 108, 800, "Onest", lh=0.98)
    tb(132, 600, 960, 240, [("данных под" + NBSP + "управлением одной платформы" + NBSP + "— на" + NBSP + "открытом "
        "lakehouse-контуре, без" + NBSP + "копий между системами и" + NBSP + "вендорского замка.", SOFT)], 23, 400, "Onest", lh=1.5)
    tb(c['x']+104, CARD1_Y+28, c['w']-134, 40, [("Lakehouse", TEXT)], 18, 700)
    tb(c['x']+104, CARD1_Y+70, c['w']-134, 110, [("StarRocks · Impala · Spark · Trino на" + NBSP + "Iceberg/S3", MUTED)], 12.5, 400, lh=1.5)
    tb(c['x']+104, CARD2_Y+28, c['w']-134, 40, [("AI Studio", TEXT)], 18, 700)
    tb(c['x']+104, CARD2_Y+70, c['w']-134, 110, [("ML · GenAI · агентный AI, MLOps/LLMOps", MUTED)], 12.5, 400, lh=1.5)
    tb(168, 984, 400, 40, [("alphyn.ai", MUTED)], 12, 400, "JetBrains Mono")

    prs.save(str(ROOT / "test_slide.pptx"))
    print("saved:", ROOT / "test_slide.pptx", "| proof:", ROOT / "slide_full.png")


if __name__ == "__main__":
    main()
